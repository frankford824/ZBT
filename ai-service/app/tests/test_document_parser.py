import json
import urllib.error
from io import BytesIO

import fitz
import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from app.pipelines.parse import document_parser
from app.pipelines.parse.document_parser import _env_int, _extract_pdf_tables, _libreoffice_convert_executable, parse_document
from app.schemas.knowledge import KnowledgeProcessRequest


def test_pdf_parser_extracts_layout_blocks_and_table_candidates() -> None:
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Item    Amount")
    page.insert_text((72, 96), "Equipment    1200")
    page.insert_text((72, 140), "Implementation plan")
    content = pdf.tobytes()
    pdf.close()

    result = parse_document(_request("layout.pdf"), content)
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert result.metadata["parser"] == "pymupdf"
    assert result.metadata["layout_block_count"] >= 1
    assert result.metadata["table_count"] >= 1
    assert result.metadata["table_block_count"] >= 1
    assert result.metadata["table_blocks"][0]["source"] == "pdf"
    assert result.metadata["table_blocks"][0]["page_start"] == 1
    assert result.metadata["page_quality_count"] == 1
    assert result.metadata["page_quality"][0]["text_char_count"] > 0
    assert result.metadata["page_quality"][0]["table_candidate_count"] >= 1
    assert result.metadata["page_quality"][0]["needs_ocr"] is False
    assert result.metadata["ocr_required"] is False
    assert "Equipment" in text


def test_table_block_md_table_preserves_empty_cells_and_escapes_pipes() -> None:
    block = document_parser._table_block(
        "ocr",
        {
            "index": 2,
            "rows": [["项目", "", "备注"], ["设备", "1200", "含|税"]],
            "cell_bboxes": [
                [[0, 0, 10, 10], [10, 0, 10, 10], [10, 0, 30, 10]],
                [[0, 10, 10, 20], [10, 10, 20, 20], [20, 10, 30, 20]],
            ],
        },
    )

    assert block["column_count"] == 3
    assert block["rows"] == [["项目", "", "备注"], ["设备", "1200", "含|税"]]
    assert block["md_table"] == "| 项目 |  | 备注 |\n| --- | --- | --- |\n| 设备 | 1200 | 含\\|税 |"
    assert block["cell_bbox_count"] == 5
    assert block["cell_bboxes"][0] == [[0.0, 0.0, 10.0, 10.0], None, [10.0, 0.0, 30.0, 10.0]]


def test_pdf_table_extraction_keeps_table_bbox() -> None:
    class FakeTable:
        bbox = (1.123, 2.0, 30.987, 40.0)
        col_count = 2
        cells = [
            (1.0, 2.0, 10.0, 8.0),
            (10.0, 2.0, 30.0, 8.0),
            (1.0, 8.0, 10.0, 20.0),
            (10.0, 8.0, 30.0, 20.0),
        ]

        def extract(self) -> list[list[str]]:
            return [["Item", "Amount"], ["Equipment", "1200"]]

    class FakeFoundTables:
        tables = [FakeTable()]

    class FakeTablePage:
        def find_tables(self) -> FakeFoundTables:
            return FakeFoundTables()

    tables, errors = _extract_pdf_tables(FakeTablePage(), 3, "")

    assert errors == []
    assert tables == [
        {
            "page": 3,
            "index": 1,
            "rows": [["Item", "Amount"], ["Equipment", "1200"]],
            "extraction": "pymupdf",
            "bbox": [1.12, 2.0, 30.99, 40.0],
            "cell_bboxes": [
                [[1.0, 2.0, 10.0, 8.0], [10.0, 2.0, 30.0, 8.0]],
                [[1.0, 8.0, 10.0, 20.0], [10.0, 8.0, 30.0, 20.0]],
            ],
        }
    ]


def test_pdf_table_extraction_error_keeps_heuristic_tables_without_sensitive_details() -> None:
    class BrokenTablePage:
        def find_tables(self) -> object:
            raise RuntimeError("secret table payload fragment")

    tables, errors = _extract_pdf_tables(
        BrokenTablePage(),
        2,
        "Item    Amount\nEquipment    1200",
    )

    assert tables == [
        {
            "page": 2,
            "index": 1,
            "rows": [["Item", "Amount"], ["Equipment", "1200"]],
            "extraction": "heuristic",
        }
    ]
    assert errors == [
        {
            "page": 2,
            "extractor": "pymupdf",
            "error_type": "RuntimeError",
            "message": "PDF 表格结构识别失败，已改用文本行识别",
        }
    ]
    assert "secret table" not in str(errors)


def test_empty_pdf_marks_ocr_required_without_claiming_success(monkeypatch) -> None:
    monkeypatch.delenv("OCR_HTTP_ENDPOINT", raising=False)
    pdf = fitz.open()
    pdf.new_page()
    content = pdf.tobytes()
    pdf.close()

    result = parse_document(_request("scan.pdf"), content)

    assert result.metadata["parser"] == "pymupdf"
    assert result.metadata["ocr_required"] is True
    assert result.metadata["ocr"]["status"] == "provider_not_configured"
    assert result.metadata["page_quality"][0]["needs_ocr"] is True
    assert result.chunks[0].metadata["needs_human_input"] is True


def test_empty_pdf_clears_ocr_required_after_successful_ocr(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.pipelines.parse.document_parser._try_http_ocr",
        lambda _payload, _content: {"status": "done", "provider": "fake_ocr", "text": "扫描件识别文本"},
    )
    pdf = fitz.open()
    pdf.new_page()
    content = pdf.tobytes()
    pdf.close()

    result = parse_document(_request("scan.pdf"), content)
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert result.metadata["ocr_required"] is False
    assert result.metadata["ocr"]["status"] == "done"
    assert "扫描件识别文本" in text


def test_mixed_pdf_runs_page_ocr_for_pages_without_text_layer(monkeypatch) -> None:
    def fake_page_ocr(payload: KnowledgeProcessRequest, content: bytes) -> dict[str, object]:
        assert payload.filename == "mixed-page-2.png"
        assert payload.content_type == "image/png"
        assert content.startswith(b"\x89PNG")
        return {
            "status": "done",
            "provider": "fake_ocr",
            "text": "第二页扫描文本",
            "metadata": {"confidence": 0.96},
        }

    monkeypatch.setenv("OCR_HTTP_ENDPOINT", "https://ocr.example.test/parse")
    monkeypatch.setattr("app.pipelines.parse.document_parser._try_http_ocr", fake_page_ocr)
    pdf = fitz.open()
    first = pdf.new_page()
    first.insert_text((72, 72), "page one selectable text")
    pdf.new_page()
    content = pdf.tobytes()
    pdf.close()

    result = parse_document(_request("mixed.pdf"), content)
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert result.metadata["ocr_required"] is False
    assert result.metadata["ocr_page_count"] == 1
    assert result.metadata["page_quality_count"] == 2
    assert result.metadata["page_quality"][1]["needs_ocr"] is True
    assert result.metadata["ocr_pages"] == [
        {"status": "done", "provider": "fake_ocr", "metadata": {"confidence": 0.96}, "page": 2}
    ]
    assert "page one selectable text" in text
    assert "第二页扫描文本" in text


def test_image_parser_uses_ocr_boundary_without_falling_back_to_plain_text(monkeypatch) -> None:
    monkeypatch.delenv("OCR_HTTP_ENDPOINT", raising=False)

    result = parse_document(_request("scan.png"), b"not-a-real-image")

    assert result.metadata["parser"] == "image-ocr"
    assert result.metadata["ocr_required"] is True
    assert result.metadata["ocr"]["status"] == "provider_not_configured"
    assert result.chunks[0].metadata["needs_human_input"] is True


def test_image_parser_clears_ocr_required_after_successful_ocr(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.pipelines.parse.document_parser._try_http_ocr",
        lambda _payload, _content: {"status": "done", "provider": "fake_ocr", "text": "图片识别文本"},
    )

    result = parse_document(_request("scan.png"), b"image-bytes")
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert result.metadata["parser"] == "image-ocr"
    assert result.metadata["ocr_required"] is False
    assert result.metadata["ocr"]["status"] == "done"
    assert "图片识别文本" in text


def test_ocr_success_response_is_normalized(monkeypatch) -> None:
    monkeypatch.setenv("OCR_HTTP_ENDPOINT", "https://ocr.example.test/parse")
    monkeypatch.setattr(
        "app.pipelines.parse.document_parser.request.urlopen",
        lambda *_args, **_kwargs: _FakeHTTPResponse(
            b"""
            {
              "pages": [
                {
                  "page": 1,
                  "text": "\xe9\xa6\x96\xe9\xa1\xb5\xe8\xaf\x86\xe5\x88\xab\xe6\x96\x87\xe6\x9c\xac",
                  "confidence": 0.92,
                  "blocks": [
                    {
                      "type": "text",
                      "text": "\xe9\xa6\x96\xe9\xa1\xb5\xe8\xaf\x86\xe5\x88\xab\xe6\x96\x87\xe6\x9c\xac",
                      "bbox": [1, 2, 30, 40],
                      "confidence": 0.91
                    }
                  ],
                  "tables": [
                    {
                      "index": 1,
                      "rows": [["\xe9\xa1\xb9", "\xe5\x80\xbc"]],
                      "confidence": 0.88
                    }
                  ]
                }
              ],
              "metadata": {"request_id": "ocr-demo"}
            }
            """
        ),
    )

    result = parse_document(_request("scan.png"), b"image-bytes")
    text = "\n".join(chunk.content for chunk in result.chunks)
    ocr = result.metadata["ocr"]

    assert result.metadata["ocr_required"] is False
    assert "首页识别文本" in text
    assert ocr["status"] == "done"
    assert ocr["confidence"] == 0.92
    assert ocr["provider_metadata"] == {"request_id": "ocr-demo"}
    assert ocr["metadata"] == {"request_id": "ocr-demo"}
    assert ocr["pages"][0]["text"] == "首页识别文本"
    assert ocr["pages"][0]["blocks"][0]["bbox"] == [1, 2, 30, 40]
    assert ocr["pages"][0]["tables"][0]["source"] == "ocr"
    assert ocr["pages"][0]["tables"][0]["confidence"] == 0.88
    assert ocr["blocks"][0]["text"] == "首页识别文本"


def test_mineru_ocr_provider_uses_specific_endpoint_key_and_mode(monkeypatch) -> None:
    captured: dict[str, object] = {}

    def fake_urlopen(req, timeout):
        captured["url"] = req.full_url
        captured["timeout"] = timeout
        captured["authorization"] = req.get_header("Authorization")
        captured["body"] = json.loads(req.data.decode("utf-8"))
        return _FakeHTTPResponse(b'{"text":"MinerU parsed text","metadata":{"request_id":"mineru-1"}}')

    monkeypatch.setenv("OCR_PROVIDER", "mineru")
    monkeypatch.setenv("MINERU_HTTP_ENDPOINT", "https://mineru.example.test/file_parse")
    monkeypatch.setenv("MINERU_API_KEY", "mineru-token")
    monkeypatch.setenv("MINERU_PARSE_MODE", "vlm")
    monkeypatch.setenv("MINERU_HTTP_TIMEOUT_S", "45")
    monkeypatch.setattr("app.pipelines.parse.document_parser.request.urlopen", fake_urlopen)

    result = parse_document(_request("scan.png"), b"image-bytes")
    ocr = result.metadata["ocr"]

    assert captured["url"] == "https://mineru.example.test/file_parse"
    assert captured["timeout"] == 45
    assert captured["authorization"] == "Bearer mineru-token"
    assert captured["body"]["provider"] == "mineru"
    assert captured["body"]["mode"] == "vlm"
    assert captured["body"]["options"]["return_markdown"] is True
    assert ocr["provider"] == "mineru"
    assert ocr["provider_profile"] == {
        "provider": "mineru",
        "endpoint_env": "MINERU_HTTP_ENDPOINT",
        "api_key_env": "MINERU_API_KEY",
        "mode": "vlm",
        "timeout_s": 45,
    }


def test_mixed_pdf_page_ocr_uses_mineru_endpoint_without_generic_endpoint(monkeypatch) -> None:
    calls: list[dict[str, object]] = []

    def fake_urlopen(req, timeout):
        calls.append({"url": req.full_url, "timeout": timeout, "body": json.loads(req.data.decode("utf-8"))})
        return _FakeHTTPResponse(b'{"data":{"text":"MinerU page OCR text","confidence":0.93}}')

    monkeypatch.setenv("OCR_PROVIDER", "mineru")
    monkeypatch.setenv("MINERU_HTTP_ENDPOINT", "https://mineru.example.test/file_parse")
    monkeypatch.delenv("OCR_HTTP_ENDPOINT", raising=False)
    monkeypatch.setattr("app.pipelines.parse.document_parser.request.urlopen", fake_urlopen)
    pdf = fitz.open()
    first = pdf.new_page()
    first.insert_text((72, 72), "selectable first page")
    pdf.new_page()
    content = pdf.tobytes()
    pdf.close()

    result = parse_document(_request("mixed.pdf"), content)
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert calls
    assert calls[0]["url"] == "https://mineru.example.test/file_parse"
    assert calls[0]["body"]["provider"] == "mineru"
    assert result.metadata["ocr_required"] is False
    assert result.metadata["ocr_page_count"] == 1
    assert result.metadata["ocr_pages"][0]["provider"] == "mineru"
    assert result.metadata["ocr_pages"][0]["provider_profile"]["endpoint_env"] == "MINERU_HTTP_ENDPOINT"
    assert "MinerU page OCR text" in text


def test_paddleocr_nested_response_normalizes_markdown_layout_and_tables(monkeypatch) -> None:
    monkeypatch.setenv("OCR_PROVIDER", "paddleocr")
    monkeypatch.setenv("PADDLEOCR_HTTP_ENDPOINT", "https://paddleocr.example.test/pp_structurev3")
    monkeypatch.setattr(
        "app.pipelines.parse.document_parser.request.urlopen",
        lambda *_args, **_kwargs: _FakeHTTPResponse(
            b"""
            {
              "data": {
                "markdown": "# \xe6\xa0\x87\xe4\xb9\xa6\xe8\xaf\x86\xe5\x88\xab\\n\xe8\xaf\x84\xe5\x88\x86\xe8\xa1\xa8",
                "layout_blocks": [
                  {
                    "page": 1,
                    "type": "title",
                    "text": "\xe6\xa0\x87\xe4\xb9\xa6\xe8\xaf\x86\xe5\x88\xab",
                    "bbox": [10, 20, 200, 60],
                    "confidence": 0.97
                  }
                ],
                "tables": [
                  {
                    "page": 1,
                    "index": 1,
                    "rows": [["\xe9\xa1\xb9", "\xe5\x88\x86\xe5\x80\xbc"], ["\xe6\x96\xb9\xe6\xa1\x88", "20"]],
                    "confidence": 0.9
                  }
                ],
                "metadata": {"request_id": "paddle-1"}
              }
            }
            """
        ),
    )

    result = parse_document(_request("scan.png"), b"image-bytes")
    text = "\n".join(chunk.content for chunk in result.chunks)
    ocr = result.metadata["ocr"]

    assert result.metadata["ocr_required"] is False
    assert result.metadata["layout_block_count"] == 1
    assert result.metadata["table_block_count"] == 1
    assert result.metadata["table_count"] == 1
    assert "标书识别" in text
    assert ocr["provider"] == "paddleocr"
    assert ocr["markdown"].startswith("# 标书识别")
    assert ocr["layout_block_count"] == 1
    assert ocr["layout_blocks"][0]["type"] == "title"
    assert ocr["layout_blocks"][0]["bbox"] == [10, 20, 200, 60]
    assert ocr["table_blocks"][0]["rows"] == [["项", "分值"], ["方案", "20"]]
    assert ocr["provider_metadata"] == {"request_id": "paddle-1"}


def test_paddleocr_provider_requires_configured_endpoint(monkeypatch) -> None:
    monkeypatch.setenv("OCR_PROVIDER", "paddleocr")
    monkeypatch.delenv("PADDLEOCR_HTTP_ENDPOINT", raising=False)
    monkeypatch.delenv("OCR_HTTP_ENDPOINT", raising=False)

    result = parse_document(_request("scan.png"), b"image-bytes")

    assert result.metadata["ocr_required"] is True
    assert result.metadata["ocr"]["status"] == "provider_not_configured"
    assert result.metadata["ocr"]["provider"] == "paddleocr"
    assert result.metadata["ocr"]["endpoint_env"] == "PADDLEOCR_HTTP_ENDPOINT"


def test_mineru_async_ocr_polls_status_url_and_normalizes_result(monkeypatch) -> None:
    calls: list[dict[str, object]] = []
    responses = [
        _FakeHTTPResponse(
            b'{"status":"processing","task_id":"mineru-task-1","status_url":"/jobs/mineru-task-1"}',
            status=202,
        ),
        _FakeHTTPResponse(b'{"status":"running","task_id":"mineru-task-1"}'),
        _FakeHTTPResponse(
            b"""
            {
              "status": "completed",
              "result": {
                "markdown": "# \xe5\xbc\x82\xe6\xad\xa5\xe8\xaf\x86\xe5\x88\xab",
                "pages": [{"page": 1, "text": "\xe5\xbc\x82\xe6\xad\xa5\xe8\xaf\x86\xe5\x88\xab\xe6\x96\x87\xe6\x9c\xac", "confidence": 0.95}],
                "tables": [{"page": 1, "index": 1, "rows": [["\xe6\x9d\xa1\xe6\xac\xbe", "\xe8\xa6\x81\xe6\xb1\x82"]]}],
                "metadata": {"request_id": "mineru-async"}
              }
            }
            """
        ),
    ]

    def fake_urlopen(req, timeout):
        calls.append(
            {
                "url": req.full_url,
                "method": req.get_method(),
                "timeout": timeout,
                "authorization": req.get_header("Authorization"),
            }
        )
        return responses.pop(0)

    monkeypatch.setenv("OCR_PROVIDER", "mineru")
    monkeypatch.setenv("MINERU_HTTP_ENDPOINT", "https://mineru.example.test/file_parse")
    monkeypatch.setenv("MINERU_API_KEY", "mineru-token")
    monkeypatch.setenv("MINERU_POLL_INTERVAL_S", "0")
    monkeypatch.setenv("MINERU_POLL_MAX_ATTEMPTS", "3")
    monkeypatch.setattr("app.pipelines.parse.document_parser.request.urlopen", fake_urlopen)

    result = parse_document(_request("scan.png"), b"image-bytes")
    text = "\n".join(chunk.content for chunk in result.chunks)
    ocr = result.metadata["ocr"]

    assert [call["method"] for call in calls] == ["POST", "GET", "GET"]
    assert calls[1]["url"] == "https://mineru.example.test/jobs/mineru-task-1"
    assert calls[1]["authorization"] == "Bearer mineru-token"
    assert result.metadata["ocr_required"] is False
    assert "异步识别文本" in text
    assert ocr["status"] == "done"
    assert ocr["table_blocks"][0]["rows"] == [["条款", "要求"]]
    assert ocr["provider_metadata"] == {
        "request_id": "mineru-async",
        "async_task_id": "mineru-task-1",
        "async_attempts": 2,
    }


def test_paddleocr_async_ocr_uses_provider_poll_endpoint_template(monkeypatch) -> None:
    calls: list[dict[str, object]] = []
    responses = [
        _FakeHTTPResponse(b'{"status":"pending","task_id":"paddle-task-1"}'),
        _FakeHTTPResponse(b'{"data":{"text":"PaddleOCR async text","confidence":0.91}}'),
    ]

    def fake_urlopen(req, timeout):
        calls.append({"url": req.full_url, "method": req.get_method(), "timeout": timeout})
        return responses.pop(0)

    monkeypatch.setenv("OCR_PROVIDER", "paddleocr")
    monkeypatch.setenv("PADDLEOCR_HTTP_ENDPOINT", "https://paddle.example.test/submit")
    monkeypatch.setenv("PADDLEOCR_POLL_ENDPOINT", "https://paddle.example.test/tasks/{task_id}/result")
    monkeypatch.setenv("PADDLEOCR_POLL_INTERVAL_S", "0")
    monkeypatch.setattr("app.pipelines.parse.document_parser.request.urlopen", fake_urlopen)

    result = parse_document(_request("scan.png"), b"image-bytes")
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert [call["method"] for call in calls] == ["POST", "GET"]
    assert calls[1]["url"] == "https://paddle.example.test/tasks/paddle-task-1/result"
    assert result.metadata["ocr"]["status"] == "done"
    assert result.metadata["ocr"]["provider_profile"]["poll_endpoint_env"] == "PADDLEOCR_POLL_ENDPOINT"
    assert "PaddleOCR async text" in text


def test_ocr_async_timeout_returns_safe_failure(monkeypatch) -> None:
    responses = [
        _FakeHTTPResponse(b'{"status":"processing","task_id":"slow-task"}', status=202),
        _FakeHTTPResponse(b'{"status":"running","debug":"secret OCR backend payload"}'),
        _FakeHTTPResponse(b'{"status":"running","debug":"secret OCR backend payload"}'),
    ]

    monkeypatch.setenv("OCR_HTTP_ENDPOINT", "https://ocr.example.test/parse")
    monkeypatch.setenv("OCR_POLL_INTERVAL_S", "0")
    monkeypatch.setenv("OCR_POLL_MAX_ATTEMPTS", "2")
    monkeypatch.setattr(
        "app.pipelines.parse.document_parser.request.urlopen",
        lambda *_args, **_kwargs: responses.pop(0),
    )

    result = parse_document(_request("scan.png"), b"image-bytes")
    ocr = result.metadata["ocr"]

    assert result.metadata["ocr_required"] is True
    assert ocr["status"] == "failed"
    assert ocr["error"] == "ocr async task timed out"
    assert ocr["task_id"] == "slow-task"
    assert ocr["attempts"] == 2
    assert "secret OCR" not in str(ocr)


def test_ocr_http_error_metadata_does_not_expose_response_body(monkeypatch) -> None:
    monkeypatch.setenv("OCR_HTTP_ENDPOINT", "https://ocr.example.test/parse")

    def raise_http_error(*_args, **_kwargs):
        raise urllib.error.HTTPError(
            "https://ocr.example.test/parse",
            502,
            "Bad Gateway",
            {},
            BytesIO(b'{"error":"secret OCR payload fragment"}'),
        )

    monkeypatch.setattr("app.pipelines.parse.document_parser.request.urlopen", raise_http_error)

    result = parse_document(_request("scan.png"), b"image-bytes")

    assert result.metadata["ocr"]["status"] == "failed"
    assert result.metadata["ocr"]["error"] == "ocr request failed"
    assert result.metadata["ocr"]["http_status"] == 502
    assert "secret OCR" not in str(result.metadata["ocr"])


@pytest.mark.parametrize(
    "endpoint",
    [
        "file:///etc/passwd",
        "https://token@ocr.example.test/parse",
        "https://ocr.example.test/parse?debug=1",
        "https://ocr.example.test/parse#fragment",
        "https://ocr.example.test\\parse",
        "https://ocr.example.test/parse\nX-Injected: yes",
        "ocr.example.test/parse",
    ],
)
def test_ocr_rejects_invalid_endpoint_without_external_request(monkeypatch, endpoint) -> None:
    monkeypatch.setenv("OCR_HTTP_ENDPOINT", endpoint)

    def fail_if_called(*_args, **_kwargs):
        raise AssertionError("invalid OCR endpoint must not be requested")

    monkeypatch.setattr("app.pipelines.parse.document_parser.request.urlopen", fail_if_called)

    result = parse_document(_request("scan.png"), b"image-bytes")

    assert result.metadata["ocr_required"] is True
    assert result.metadata["ocr"]["status"] == "failed"
    assert result.metadata["ocr"]["error"] == "ocr endpoint invalid"
    assert "token@" not in str(result.metadata["ocr"])


def test_ocr_rejects_invalid_api_key_without_external_request(monkeypatch) -> None:
    monkeypatch.setenv("OCR_HTTP_ENDPOINT", "https://ocr.example.test/parse")
    monkeypatch.setenv("OCR_API_KEY", "secret\r\nX-Injected: yes")

    def fail_if_called(*_args, **_kwargs):
        raise AssertionError("invalid OCR API key must not be requested")

    monkeypatch.setattr("app.pipelines.parse.document_parser.request.urlopen", fail_if_called)

    result = parse_document(_request("scan.png"), b"image-bytes")

    assert result.metadata["ocr_required"] is True
    assert result.metadata["ocr"]["status"] == "failed"
    assert result.metadata["ocr"]["error"] == "ocr api key invalid"
    assert "secret" not in str(result.metadata["ocr"])


def test_ocr_skips_oversized_content_without_external_request(monkeypatch) -> None:
    monkeypatch.setenv("OCR_HTTP_ENDPOINT", "https://ocr.example.test/parse")
    monkeypatch.setenv("OCR_MAX_BYTES", "4")

    def fail_if_called(*_args, **_kwargs):
        raise AssertionError("oversized OCR content must not be sent")

    monkeypatch.setattr("app.pipelines.parse.document_parser.request.urlopen", fail_if_called)

    result = parse_document(_request("scan.png"), b"image-bytes")

    assert result.metadata["ocr_required"] is True
    assert result.metadata["ocr"]["status"] == "skipped_too_large"
    assert result.metadata["ocr"]["size_bytes"] == len(b"image-bytes")
    assert result.metadata["ocr"]["max_bytes"] == 4


def test_ocr_rejects_oversized_response_without_exposing_body(monkeypatch) -> None:
    monkeypatch.setenv("OCR_HTTP_ENDPOINT", "https://ocr.example.test/parse")
    monkeypatch.setenv("OCR_MAX_RESPONSE_BYTES", "1024")
    oversized_body = b'{"text":"' + (b"secret OCR response body" * 80) + b'"}'

    monkeypatch.setattr(
        "app.pipelines.parse.document_parser.request.urlopen",
        lambda *_args, **_kwargs: _FakeHTTPResponse(oversized_body),
    )

    result = parse_document(_request("scan.png"), b"image-bytes")

    assert result.metadata["ocr"]["status"] == "failed"
    assert result.metadata["ocr"]["error"] == "ocr response too large"
    assert "secret OCR" not in str(result.metadata["ocr"])


def test_ocr_rejects_non_object_response_without_exposing_body(monkeypatch) -> None:
    monkeypatch.setenv("OCR_HTTP_ENDPOINT", "https://ocr.example.test/parse")

    monkeypatch.setattr(
        "app.pipelines.parse.document_parser.request.urlopen",
        lambda *_args, **_kwargs: _FakeHTTPResponse(b'["secret OCR response body"]'),
    )

    result = parse_document(_request("scan.png"), b"image-bytes")

    assert result.metadata["ocr"]["status"] == "failed"
    assert result.metadata["ocr"]["error"] == "ocr response invalid"
    assert "secret OCR" not in str(result.metadata["ocr"])


def test_legacy_office_marks_human_input_when_converter_missing(monkeypatch) -> None:
    monkeypatch.delenv("LIBREOFFICE_BIN", raising=False)
    monkeypatch.delenv("LIBREOFFICE_PATH", raising=False)
    monkeypatch.setattr("app.pipelines.parse.document_parser.shutil.which", lambda _name: None)

    result = parse_document(_request("legacy.doc"), b"legacy-binary")

    assert result.metadata["parser"] == "legacy-office-unconverted"
    assert result.metadata["legacy_conversion"]["status"] == "converter_not_configured"
    assert result.metadata["needs_human_input"] is True
    assert result.chunks[0].metadata["needs_human_input"] is True


def test_legacy_office_conversion_failure_does_not_expose_converter_output(monkeypatch) -> None:
    class FailedConversion:
        returncode = 1
        stderr = "secret tender body fragment"
        stdout = "C:\\sensitive\\legacy.doc"

    monkeypatch.setattr("app.pipelines.parse.document_parser._libreoffice_convert_executable", lambda: "/usr/bin/soffice")
    monkeypatch.setattr("app.pipelines.parse.document_parser.subprocess.run", lambda *_args, **_kwargs: FailedConversion())

    result = parse_document(_request("legacy.doc"), b"legacy-binary")
    conversion = result.metadata["legacy_conversion"]

    assert conversion == {
        "status": "failed",
        "target_suffix": ".docx",
        "error": "conversion_failed",
        "return_code": 1,
    }
    assert "secret tender" not in str(conversion)
    assert "sensitive" not in str(conversion)
    assert result.metadata["needs_human_input"] is True


def test_legacy_office_converter_uses_libreoffice_path(monkeypatch) -> None:
    monkeypatch.delenv("LIBREOFFICE_BIN", raising=False)
    monkeypatch.setenv("LIBREOFFICE_PATH", "/opt/libreoffice/program/soffice")
    monkeypatch.setattr("app.pipelines.parse.document_parser.shutil.which", lambda _name: None)

    assert _libreoffice_convert_executable() == "/opt/libreoffice/program/soffice"


def test_docx_parser_includes_paragraphs_and_tables() -> None:
    document = Document()
    document.add_paragraph("项目总体方案")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "工期"
    table.cell(0, 1).text = "90天"
    content = BytesIO()
    document.save(content)

    result = parse_document(_request("plan.docx"), content.getvalue())
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert result.metadata["parser"] == "python-docx"
    assert "项目总体方案" in text
    assert "工期 | 90天" in text
    assert result.metadata["docx_paragraph_count"] == 1
    assert result.metadata["docx_table_count"] == 1
    assert result.metadata["table_count"] == 1
    assert result.metadata["table_blocks"] == [
        {
            "source": "docx",
            "index": 1,
            "rows": [["工期", "90天"]],
            "row_count": 1,
            "column_count": 2,
            "md_table": "| 工期 | 90天 |\n| --- | --- |",
            "extraction": "python-docx",
        }
    ]
    assert result.metadata["truncated_after_parse_limit"] is False


def test_docx_parser_marks_configured_parse_limits(monkeypatch) -> None:
    monkeypatch.setenv("KNOWLEDGE_PARSE_MAX_DOCX_PARAGRAPHS", "2")
    monkeypatch.setenv("KNOWLEDGE_PARSE_MAX_DOCX_TABLE_ROWS", "1")
    document = Document()
    document.add_paragraph("第一段方案")
    document.add_paragraph("第二段方案")
    document.add_paragraph("第三段方案")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "首行"
    table.cell(0, 1).text = "保留"
    table.cell(1, 0).text = "次行"
    table.cell(1, 1).text = "截断"
    content = BytesIO()
    document.save(content)

    result = parse_document(_request("plan.docx"), content.getvalue())
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert "第一段方案" in text
    assert "第二段方案" in text
    assert "第三段方案" not in text
    assert "首行 | 保留" in text
    assert "次行 | 截断" not in text
    assert result.metadata["docx_paragraph_limit"] == 2
    assert result.metadata["docx_paragraph_count"] == 3
    assert result.metadata["docx_parsed_paragraph_count"] == 2
    assert result.metadata["docx_table_row_limit"] == 1
    assert result.metadata["docx_table_row_count"] == 2
    assert result.metadata["docx_parsed_table_row_count"] == 1
    assert result.metadata["truncated_after_parse_limit"] is True


def test_xlsx_parser_extracts_sheet_rows() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "报价"
    sheet.append(["科目", "金额"])
    sheet.append(["设备", 1200])
    content = BytesIO()
    workbook.save(content)

    result = parse_document(_request("quote.xlsx"), content.getvalue())
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert result.metadata["parser"] == "openpyxl"
    assert "[Sheet] 报价" in text
    assert "设备 | 1200" in text
    assert result.metadata["table_blocks"][0]["source"] == "xlsx"
    assert result.metadata["table_blocks"][0]["sheet"] == "报价"
    assert result.metadata["table_blocks"][0]["rows"] == [["科目", "金额"], ["设备", "1200"]]
    assert result.metadata["table_blocks"][0]["md_table"] == "| 科目 | 金额 |\n| --- | --- |\n| 设备 | 1200 |"


def test_xlsx_parser_preserves_uncached_formula_text() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "报价"
    sheet.append(["科目", "金额"])
    sheet.append(["设备", 1200])
    sheet.append(["人工", 800])
    sheet.append(["合计", "=SUM(B2:B3)"])
    content = BytesIO()
    workbook.save(content)

    result = parse_document(_request("quote.xlsx"), content.getvalue())
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert "设备 | 1200" in text
    assert "合计 | =SUM(B2:B3)" in text


def test_xlsx_parser_closes_primary_workbook_when_formula_open_fails(monkeypatch) -> None:
    class FakeWorkbook:
        worksheets: list[object] = []

        def __init__(self) -> None:
            self.closed = False

        def close(self) -> None:
            self.closed = True

    workbook = FakeWorkbook()
    calls = 0

    def fake_load_workbook(*_args, **_kwargs):
        nonlocal calls
        calls += 1
        if calls == 1:
            return workbook
        raise RuntimeError("formula workbook open failed")

    monkeypatch.setattr(document_parser, "load_workbook", fake_load_workbook)

    with pytest.raises(RuntimeError, match="formula workbook open failed"):
        document_parser._parse_xlsx(b"xlsx-bytes")

    assert workbook.closed is True


def test_xlsx_parser_stops_at_configured_row_limit(monkeypatch) -> None:
    monkeypatch.setenv("KNOWLEDGE_PARSE_MAX_XLSX_ROWS_PER_SHEET", "2")
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "报价"
    sheet.append(["科目", "金额"])
    sheet.append(["设备", 1200])
    sheet.append(["人工", 800])
    content = BytesIO()
    workbook.save(content)

    result = parse_document(_request("quote.xlsx"), content.getvalue())
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert "设备 | 1200" in text
    assert "人工 | 800" not in text
    assert result.metadata["xlsx_row_limit_per_sheet"] == 2
    assert result.metadata["truncated_after_parse_limit"] is True


def test_pptx_parser_extracts_slide_text() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "实施路线"
    textbox = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    textbox.text = "里程碑计划"
    content = BytesIO()
    presentation.save(content)

    result = parse_document(_request("deck.pptx"), content.getvalue())
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert result.metadata["parser"] == "python-pptx"
    assert "实施路线" in text
    assert "里程碑计划" in text


def test_pptx_parser_extracts_table_blocks() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    table_shape = slide.shapes.add_table(1, 2, 0, 0, 1000000, 1000000)
    table_shape.table.cell(0, 0).text = "节点"
    table_shape.table.cell(0, 1).text = "完成"
    content = BytesIO()
    presentation.save(content)

    result = parse_document(_request("deck.pptx"), content.getvalue())
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert "节点 | 完成" in text
    assert result.metadata["table_blocks"] == [
        {
            "source": "pptx",
            "index": 1,
            "rows": [["节点", "完成"]],
            "row_count": 1,
            "column_count": 2,
            "md_table": "| 节点 | 完成 |\n| --- | --- |",
            "slide": 1,
            "extraction": "python-pptx",
        }
    ]


def test_pptx_parser_stops_at_configured_slide_limit(monkeypatch) -> None:
    monkeypatch.setenv("KNOWLEDGE_PARSE_MAX_PPTX_SLIDES", "1")
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[5])
    first.shapes.title.text = "第一页方案"
    second = presentation.slides.add_slide(presentation.slide_layouts[5])
    second.shapes.title.text = "第二页方案"
    content = BytesIO()
    presentation.save(content)

    result = parse_document(_request("deck.pptx"), content.getvalue())
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert "第一页方案" in text
    assert "第二页方案" not in text
    assert result.metadata["pptx_slide_limit"] == 1
    assert result.metadata["pptx_parsed_slide_count"] == 1
    assert result.metadata["truncated_after_parse_limit"] is True


def test_pdf_parser_stops_at_configured_page_limit(monkeypatch) -> None:
    monkeypatch.setenv("KNOWLEDGE_PARSE_MAX_PDF_PAGES", "1")
    pdf = fitz.open()
    first = pdf.new_page()
    first.insert_text((72, 72), "first page content")
    second = pdf.new_page()
    second.insert_text((72, 72), "second page content")
    content = pdf.tobytes()
    pdf.close()

    result = parse_document(_request("multi.pdf"), content)
    text = "\n".join(chunk.content for chunk in result.chunks)

    assert "first page content" in text
    assert "second page content" not in text
    assert result.metadata["page_count"] == 2
    assert result.metadata["parsed_page_count"] == 1
    assert result.metadata["truncated_after_page_limit"] is True


def test_parse_document_marks_chunk_limit_truncation(monkeypatch) -> None:
    monkeypatch.setenv("KNOWLEDGE_PARSE_MAX_CHUNKS", "2")
    content = ("\n".join(f"段落 {index} " + "内容" * 40 for index in range(80))).encode()

    result = parse_document(_request("large.txt"), content)

    assert len(result.chunks) == 2
    assert result.metadata["chunk_count"] == 2
    assert result.metadata["chunk_limit"] == 2
    assert result.metadata["truncated_after_chunk_limit"] is True
    assert result.chunks[-1].metadata["truncated_after_chunk_limit"] is True


def test_timeout_env_parsing_falls_back_for_invalid_values(monkeypatch) -> None:
    monkeypatch.setenv("OCR_HTTP_TIMEOUT_S", "bad")
    assert _env_int("OCR_HTTP_TIMEOUT_S", 120) == 120

    monkeypatch.setenv("OCR_HTTP_TIMEOUT_S", "0")
    assert _env_int("OCR_HTTP_TIMEOUT_S", 120) == 120

    monkeypatch.setenv("OCR_HTTP_TIMEOUT_S", "30")
    assert _env_int("OCR_HTTP_TIMEOUT_S", 120) == 30


def _request(filename: str) -> KnowledgeProcessRequest:
    return KnowledgeProcessRequest(
        tenant_id="tenant-demo",
        document_id="doc-demo",
        file_id="file-demo",
        object_key=f"demo/{filename}",
        filename=filename,
        content_type="application/octet-stream",
    )


class _FakeHTTPResponse:
    def __init__(self, body: bytes, status: int = 200) -> None:
        self._body = BytesIO(body)
        self.status = status
        self.code = status

    def __enter__(self) -> "_FakeHTTPResponse":
        return self

    def __exit__(self, *_args: object) -> None:
        return None

    def read(self, size: int = -1) -> bytes:
        return self._body.read(size)
