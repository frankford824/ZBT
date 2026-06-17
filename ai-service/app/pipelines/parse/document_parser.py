from __future__ import annotations

import base64
import json
import os
import re
import shutil
import subprocess
from dataclasses import dataclass
from itertools import zip_longest
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory
from urllib import error, parse, request

import fitz
from openpyxl import load_workbook
from docx import Document as DocxDocument
from pptx import Presentation

from app.schemas.knowledge import KnowledgeChunk, KnowledgeProcessRequest, KnowledgeProcessResult


IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
LEGACY_OFFICE_TARGETS = {
    ".doc": ".docx",
    ".xls": ".xlsx",
    ".ppt": ".pptx",
}
DEFAULT_OCR_RESPONSE_MAX_BYTES = 8 * 1024 * 1024
SUPPORTED_OCR_PROVIDERS = {"http_ocr", "http", "mineru", "paddleocr"}


class OCRResponseTooLargeError(Exception):
    pass


@dataclass(frozen=True)
class OCRProviderConfig:
    provider: str
    endpoint: str
    endpoint_env: str
    api_key: str
    api_key_env: str
    timeout_s: int
    mode: str


def _env_int(name: str, default: int, minimum: int = 1) -> int:
    raw = os.getenv(name, "").strip()
    if not raw:
        return default
    try:
        value = int(raw)
    except ValueError:
        return default
    return value if value >= minimum else default


def parse_document(payload: KnowledgeProcessRequest, content: bytes) -> KnowledgeProcessResult:
    suffix = Path(payload.filename).suffix.lower()
    content_type = payload.content_type.lower()
    metadata: dict[str, object] = {
        "content_type": payload.content_type,
        "object_key": payload.object_key,
    }
    if "pdf" in content_type or suffix == ".pdf":
        text, page_count, pdf_metadata = _parse_pdf(payload, content)
        metadata.update(pdf_metadata)
        if not text.strip():
            ocr_result = _try_http_ocr(payload, content)
            if ocr_result["status"] == "done":
                text = str(ocr_result.get("text") or "")
                ocr_result = {key: value for key, value in ocr_result.items() if key != "text"}
            metadata["ocr"] = ocr_result
            metadata["ocr_required"] = not bool(text.strip())
        parser = "pymupdf"
    elif suffix in LEGACY_OFFICE_TARGETS:
        text, page_count, legacy_metadata = _parse_legacy_office(payload, content, suffix)
        metadata.update(legacy_metadata)
        parser = str(legacy_metadata.get("parser") or "libreoffice-conversion")
    elif "word" in content_type or suffix == ".docx":
        text, parse_metadata = _parse_docx(content)
        page_count = None
        metadata.update(parse_metadata)
        parser = "python-docx"
    elif "spreadsheet" in content_type or suffix in {".xlsx", ".xlsm"}:
        text, parse_metadata = _parse_xlsx(content)
        page_count = None
        metadata.update(parse_metadata)
        parser = "openpyxl"
    elif "presentation" in content_type or suffix in {".pptx", ".pptm"}:
        text, parse_metadata = _parse_pptx(content)
        page_count = None
        metadata.update(parse_metadata)
        parser = "python-pptx"
    elif content_type.startswith("image/") or suffix in IMAGE_SUFFIXES:
        ocr_result = _try_http_ocr(payload, content)
        text = str(ocr_result.get("text") or "") if ocr_result["status"] == "done" else ""
        metadata["ocr"] = {key: value for key, value in ocr_result.items() if key != "text"}
        metadata["ocr_required"] = not bool(text.strip())
        page_count = None
        parser = "image-ocr"
    else:
        text, page_count = _parse_text(content), None
        parser = "plain-text"

    chunk_limit = _env_int("KNOWLEDGE_PARSE_MAX_CHUNKS", 300)
    chunks = _chunk_text(text, payload.filename, page_count, chunk_limit)
    summary = _summary(text, payload.filename)
    metadata.update(
        {
            "parser": parser,
            "page_count": page_count,
            "chunk_count": len(chunks),
            "chunk_limit": chunk_limit,
            "truncated_after_chunk_limit": any(
                bool(chunk.metadata.get("truncated_after_chunk_limit")) for chunk in chunks
            ),
        }
    )
    return KnowledgeProcessResult(
        processed_title=Path(payload.filename).stem or payload.filename,
        summary=summary,
        chunks=chunks,
        metadata=metadata,
    )


def _parse_text(content: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "gb18030"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _parse_pdf(payload: KnowledgeProcessRequest, content: bytes) -> tuple[str, int, dict[str, object]]:
    doc = fitz.open(stream=content, filetype="pdf")
    try:
        page_limit = _env_int("KNOWLEDGE_PARSE_MAX_PDF_PAGES", 300)
        parsed_page_count = min(doc.page_count, page_limit)
        page_texts: list[str] = []
        layout_blocks: list[dict[str, object]] = []
        tables: list[dict[str, object]] = []
        page_quality: list[dict[str, object]] = []
        table_errors: list[dict[str, object]] = []
        for page_index in range(1, parsed_page_count + 1):
            page = doc[page_index - 1]
            page_text = page.get_text("text").strip()
            raw_layout = page.get_text("dict")
            page_texts.append(page_text)
            layout_blocks.extend(_extract_pdf_layout_blocks(raw_layout, page_index))
            page_tables, page_table_errors = _extract_pdf_tables(page, page_index, page_text)
            tables.extend(page_tables)
            table_errors.extend(page_table_errors)
            page_quality.append(_pdf_page_quality(page, raw_layout, page_text, page_tables, page_index))
        pages: list[str] = []
        page_ocr_results: list[dict[str, object]] = []
        has_text_layer = any(text.strip() for text in page_texts)
        for page_index, page_text in enumerate(page_texts, start=1):
            if page_text:
                pages.append(f"[Page {page_index}]\n{page_text}")
                continue
            if not has_text_layer:
                continue
            ocr_result = _try_pdf_page_ocr(payload, doc[page_index - 1], page_index)
            ocr_text = str(ocr_result.get("text") or "")
            ocr_metadata = {key: value for key, value in ocr_result.items() if key != "text"}
            ocr_metadata["page"] = page_index
            page_ocr_results.append(ocr_metadata)
            if ocr_result.get("status") == "done" and ocr_text.strip():
                pages.append(f"[Page {page_index} OCR]\n{ocr_text.strip()}")
        text = "\n\n".join(pages)
        metadata = {
            "layout_blocks": layout_blocks[:200],
            "layout_block_count": len(layout_blocks),
            "tables": tables[:50],
            "table_count": len(tables),
            "table_blocks": [_table_block("pdf", table) for table in tables[:50]],
            "table_block_count": len(tables),
            "page_quality": page_quality[:200],
            "page_quality_count": len(page_quality),
            "table_extraction_errors": table_errors[:20],
            "table_extraction_error_count": len(table_errors),
            "ocr_pages": page_ocr_results[:50],
            "ocr_page_count": len(page_ocr_results),
            "ocr_required": _pdf_ocr_required(text, page_ocr_results),
            "parsed_page_count": parsed_page_count,
            "page_limit": page_limit,
            "truncated_after_page_limit": doc.page_count > parsed_page_count,
        }
        return text, doc.page_count, metadata
    finally:
        doc.close()


def _pdf_ocr_required(text: str, page_ocr_results: list[dict[str, object]]) -> bool:
    if not text.strip():
        return True
    return any(result.get("status") != "done" for result in page_ocr_results)


def _try_pdf_page_ocr(payload: KnowledgeProcessRequest, page: fitz.Page, page_index: int) -> dict[str, object]:
    if not os.getenv("OCR_HTTP_ENDPOINT", "").strip():
        return {"status": "provider_not_configured", "provider": _ocr_provider_name()}
    try:
        pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        page_image = pixmap.tobytes("png")
    except Exception as exc:
        return {
            "status": "failed",
            "provider": _ocr_provider_name(),
            "error": "pdf page render failed",
            "error_type": type(exc).__name__,
        }
    page_payload = payload.model_copy(
        update={
            "filename": f"{Path(payload.filename).stem or 'page'}-page-{page_index}.png",
            "content_type": "image/png",
        }
    )
    return _try_http_ocr(page_payload, page_image)


def _extract_pdf_layout_blocks(raw: dict[str, object], page_index: int) -> list[dict[str, object]]:
    blocks: list[dict[str, object]] = []
    for block in raw.get("blocks", []):
        if block.get("type") != 0:
            continue
        lines: list[str] = []
        for line in block.get("lines", []):
            line_text = "".join(span.get("text", "") for span in line.get("spans", [])).strip()
            if line_text:
                lines.append(line_text)
        text = "\n".join(lines).strip()
        if not text:
            continue
        blocks.append(
            {
                "page": page_index,
                "bbox": [round(float(value), 2) for value in block.get("bbox", [])],
                "text": text[:1000],
            }
        )
    return blocks


def _pdf_page_quality(
    page: fitz.Page,
    raw: dict[str, object],
    page_text: str,
    tables: list[dict[str, object]],
    page_index: int,
) -> dict[str, object]:
    page_area = max(float(page.rect.width * page.rect.height), 1.0)
    text_area = 0.0
    image_area = 0.0
    image_block_count = 0
    text_block_count = 0
    for block in raw.get("blocks", []):
        bbox = block.get("bbox")
        area = _bbox_area(bbox)
        if block.get("type") == 0:
            text_block_count += 1
            text_area += area
        elif block.get("type") == 1:
            image_block_count += 1
            image_area += area
    text_char_count = len(page_text.strip())
    text_density = round(text_char_count / page_area * 1000, 4)
    text_coverage = round(min(text_area / page_area, 1.0), 4)
    image_area_ratio = round(min(image_area / page_area, 1.0), 4)
    needs_ocr = text_char_count < 20
    return {
        "page": page_index,
        "text_char_count": text_char_count,
        "text_block_count": text_block_count,
        "text_density": text_density,
        "text_coverage": text_coverage,
        "image_block_count": image_block_count,
        "image_area_ratio": image_area_ratio,
        "table_candidate_count": len(tables),
        "needs_ocr": needs_ocr,
    }


def _bbox_area(value: object) -> float:
    if not isinstance(value, (list, tuple)) or len(value) != 4:
        return 0.0
    try:
        x0, y0, x1, y1 = (float(item) for item in value)
    except (TypeError, ValueError):
        return 0.0
    return max(x1 - x0, 0.0) * max(y1 - y0, 0.0)


def _extract_pdf_tables(
    page: fitz.Page,
    page_index: int,
    page_text: str,
) -> tuple[list[dict[str, object]], list[dict[str, object]]]:
    tables: list[dict[str, object]] = []
    errors: list[dict[str, object]] = []
    if hasattr(page, "find_tables"):
        try:
            max_rows = _env_int("KNOWLEDGE_PARSE_MAX_TABLE_ROWS", 200)
            found = page.find_tables()
            for table_index, table in enumerate(found.tables, start=1):
                rows = table.extract()
                cleaned_rows: list[list[str]] = []
                for row in rows[:max_rows]:
                    values = [str(cell or "").strip() for cell in row]
                    if any(values):
                        cleaned_rows.append(values)
                if cleaned_rows:
                    table_payload: dict[str, object] = {"page": page_index, "index": table_index, "rows": cleaned_rows}
                    if len(rows) > max_rows:
                        table_payload["truncated_after_row_limit"] = True
                        table_payload["row_limit"] = max_rows
                    tables.append(table_payload)
        except Exception as exc:
            errors.append(
                {
                    "page": page_index,
                    "extractor": "pymupdf",
                    "error_type": type(exc).__name__,
                    "message": "PDF 表格结构识别失败，已改用文本行识别",
                }
            )
    if tables:
        return tables, errors

    heuristic_rows: list[list[str]] = []
    max_rows = _env_int("KNOWLEDGE_PARSE_MAX_TABLE_ROWS", 200)
    for line in page_text.splitlines():
        cells = [cell.strip() for cell in re.split(r"\t+|\s{2,}|\s*\|\s*", line.strip()) if cell.strip()]
        if len(cells) >= 2:
            heuristic_rows.append(cells)
            if len(heuristic_rows) >= max_rows:
                break
    if len(heuristic_rows) >= 2:
        table_payload = {"page": page_index, "index": 1, "rows": heuristic_rows, "extraction": "heuristic"}
        if len(heuristic_rows) >= max_rows:
            table_payload["truncated_after_row_limit"] = True
            table_payload["row_limit"] = max_rows
        tables.append(table_payload)
    return tables, errors


def _table_block(source: str, table: dict[str, object]) -> dict[str, object]:
    rows = table.get("rows")
    normalized_rows = rows if isinstance(rows, list) else []
    block: dict[str, object] = {
        "source": source,
        "index": _metadata_int(table.get("index"), 1),
        "rows": normalized_rows,
        "row_count": len(normalized_rows),
    }
    page = _metadata_int(table.get("page"), 0)
    if page:
        block["page_start"] = page
        block["page_end"] = page
    sheet = str(table.get("sheet") or "").strip()
    if sheet:
        block["sheet"] = sheet
    slide = _metadata_int(table.get("slide"), 0)
    if slide:
        block["slide"] = slide
    extraction = str(table.get("extraction") or "").strip()
    if extraction:
        block["extraction"] = extraction
    if table.get("truncated_after_row_limit"):
        block["truncated_after_row_limit"] = True
        block["row_limit"] = _metadata_int(table.get("row_limit"), 0)
    bbox = table.get("bbox")
    if isinstance(bbox, list):
        block["bbox"] = bbox
    confidence = _metadata_float(table.get("confidence"))
    if confidence is not None:
        block["confidence"] = confidence
    return block


def _metadata_int(value: object, default: int) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _metadata_float(value: object) -> float | None:
    if value in (None, ""):
        return None
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _normalize_ocr_result(provider: str, result: dict[str, object]) -> dict[str, object]:
    provider_metadata = result.get("provider_metadata")
    if not isinstance(provider_metadata, dict):
        provider_metadata = result.get("metadata") if isinstance(result.get("metadata"), dict) else {}
    pages = _normalize_ocr_pages(result.get("pages"))
    text = str(result.get("text") or "").strip()
    if not text and pages:
        text = "\n\n".join(str(page.get("text") or "") for page in pages).strip()
    blocks = _normalize_ocr_blocks(result.get("blocks"), pages)
    raw_tables = result.get("tables")
    tables = [
        _table_block("ocr", table)
        for table in raw_tables
        if isinstance(table, dict)
    ] if isinstance(raw_tables, list) else []
    confidence = _ocr_confidence(result, pages)
    return {
        "status": "done" if text.strip() else "empty_result",
        "provider": provider,
        "text": text,
        "pages": pages,
        "blocks": blocks,
        "tables": tables,
        "table_blocks": tables,
        "confidence": confidence,
        "provider_metadata": provider_metadata,
        "metadata": provider_metadata,
    }


def _normalize_ocr_pages(raw_pages: object) -> list[dict[str, object]]:
    if not isinstance(raw_pages, list):
        return []
    pages: list[dict[str, object]] = []
    for index, item in enumerate(raw_pages, start=1):
        if not isinstance(item, dict):
            continue
        text = str(item.get("text") or "")
        page: dict[str, object] = {
            "page": _metadata_int(item.get("page") or item.get("page_index"), index),
            "text": text[:2000],
            "text_char_count": len(text),
        }
        if len(text) > 2000:
            page["text_truncated"] = True
        confidence = _metadata_float(item.get("confidence"))
        if confidence is not None:
            page["confidence"] = confidence
        blocks = _normalize_ocr_blocks(item.get("blocks"), [])
        if blocks:
            page["blocks"] = blocks[:100]
            page["block_count"] = len(blocks)
        raw_tables = item.get("tables")
        if isinstance(raw_tables, list):
            tables = [
                _table_block("ocr", {**table, "page": page["page"]})
                for table in raw_tables
                if isinstance(table, dict)
            ]
            if tables:
                page["tables"] = tables[:30]
                page["table_count"] = len(tables)
        pages.append(page)
    return pages


def _normalize_ocr_blocks(raw_blocks: object, pages: list[dict[str, object]]) -> list[dict[str, object]]:
    blocks: list[dict[str, object]] = []
    if isinstance(raw_blocks, list):
        for item in raw_blocks[:500]:
            if not isinstance(item, dict):
                continue
            text = str(item.get("text") or "").strip()
            block: dict[str, object] = {
                "text": text[:1000],
                "text_char_count": len(text),
            }
            page = _metadata_int(item.get("page") or item.get("page_index"), 0)
            if page:
                block["page"] = page
            bbox = item.get("bbox")
            if isinstance(bbox, list):
                block["bbox"] = bbox
            block_type = str(item.get("type") or "").strip()
            if block_type:
                block["type"] = block_type
            confidence = _metadata_float(item.get("confidence"))
            if confidence is not None:
                block["confidence"] = confidence
            blocks.append(block)
    if not blocks:
        for page in pages:
            text = str(page.get("text") or "").strip()
            if text:
                blocks.append(
                    {
                        "page": page.get("page"),
                        "type": "text",
                        "text": text[:1000],
                        "text_char_count": int(page.get("text_char_count") or len(text)),
                    }
                )
    return blocks


def _ocr_confidence(result: dict[str, object], pages: list[dict[str, object]]) -> float | None:
    confidence = _metadata_float(result.get("confidence"))
    if confidence is not None:
        return confidence
    metadata = result.get("metadata")
    if isinstance(metadata, dict):
        confidence = _metadata_float(metadata.get("confidence"))
        if confidence is not None:
            return confidence
    page_confidences = [
        value
        for value in (_metadata_float(page.get("confidence")) for page in pages)
        if value is not None
    ]
    if not page_confidences:
        return None
    return round(sum(page_confidences) / len(page_confidences), 4)


def _try_http_ocr(payload: KnowledgeProcessRequest, content: bytes) -> dict[str, object]:
    config = _ocr_provider_config()
    if not config.endpoint:
        return {
            "status": "provider_not_configured",
            "provider": config.provider,
            "endpoint_env": config.endpoint_env,
            "supported_providers": sorted(SUPPORTED_OCR_PROVIDERS),
        }
    try:
        endpoint = _safe_ocr_endpoint(config.endpoint)
    except RuntimeError:
        return {"status": "failed", "provider": config.provider, "error": "ocr endpoint invalid"}
    max_bytes = _env_int("OCR_MAX_BYTES", 20 * 1024 * 1024)
    if len(content) > max_bytes:
        return {
            "status": "skipped_too_large",
            "provider": config.provider,
            "size_bytes": len(content),
            "max_bytes": max_bytes,
        }
    body = json.dumps(
        {
            "filename": payload.filename,
            "content_type": payload.content_type,
            "provider": config.provider,
            "mode": config.mode,
            "options": _ocr_provider_options(config),
            "content_base64": base64.b64encode(content).decode("ascii"),
        },
        ensure_ascii=False,
    ).encode("utf-8")
    headers = {"Content-Type": "application/json"}
    if config.api_key:
        try:
            headers["Authorization"] = f"Bearer {_safe_ocr_header_value(config.api_key)}"
        except RuntimeError:
            return {"status": "failed", "provider": config.provider, "error": "ocr api key invalid"}
    req = request.Request(
        endpoint,
        data=body,
        method="POST",
        headers=headers,
    )
    try:
        with request.urlopen(req, timeout=config.timeout_s) as response:
            max_response_bytes = _env_int(
                "OCR_MAX_RESPONSE_BYTES",
                DEFAULT_OCR_RESPONSE_MAX_BYTES,
                minimum=1024,
            )
            result = json.loads(_read_limited_response(response, max_response_bytes).decode("utf-8"))
        if not isinstance(result, dict):
            return {
                "status": "failed",
                "provider": config.provider,
                "error": "ocr response invalid",
            }
        normalized = _normalize_ocr_result(config.provider, result)
        normalized["provider_profile"] = _ocr_provider_profile(config)
        return normalized
    except error.HTTPError as exc:
        return {
            "status": "failed",
            "provider": config.provider,
            "error": "ocr request failed",
            "http_status": exc.code,
        }
    except OCRResponseTooLargeError:
        return {
            "status": "failed",
            "provider": config.provider,
            "error": "ocr response too large",
        }
    except Exception:
        return {"status": "failed", "provider": config.provider, "error": "ocr request failed"}


def _safe_ocr_endpoint(value: str) -> str:
    if _contains_unsafe_url_character(value):
        raise RuntimeError("OCR_HTTP_ENDPOINT is invalid")
    parsed = parse.urlparse(value)
    if (
        parsed.scheme not in {"http", "https"}
        or not parsed.netloc
        or not parsed.hostname
        or parsed.username is not None
        or parsed.password is not None
        or parsed.params
        or parsed.query
        or parsed.fragment
    ):
        raise RuntimeError("OCR_HTTP_ENDPOINT must be an absolute HTTP(S) URL")
    return parse.urlunparse((parsed.scheme, parsed.netloc, parsed.path or "/", "", "", ""))


def _contains_unsafe_url_character(value: str) -> bool:
    return "\\" in value or any(ord(char) <= 0x20 or ord(char) == 0x7F for char in value)


def _safe_ocr_header_value(value: str) -> str:
    cleaned = value.strip()
    if not cleaned or any(char in cleaned for char in "\r\n"):
        raise RuntimeError("OCR header value is invalid")
    return cleaned


def _read_limited_response(response: object, max_bytes: int) -> bytes:
    content = response.read(max_bytes + 1)
    if len(content) > max_bytes:
        raise OCRResponseTooLargeError
    return content


def _ocr_provider_name() -> str:
    provider = os.getenv("OCR_PROVIDER", "http_ocr").strip().lower() or "http_ocr"
    return provider if provider in SUPPORTED_OCR_PROVIDERS else "http_ocr"


def _ocr_provider_config() -> OCRProviderConfig:
    provider = _ocr_provider_name()
    endpoint_env = _ocr_endpoint_env(provider)
    api_key_env = _ocr_api_key_env(provider)
    return OCRProviderConfig(
        provider=provider,
        endpoint=os.getenv(endpoint_env, "").strip() or os.getenv("OCR_HTTP_ENDPOINT", "").strip(),
        endpoint_env=endpoint_env,
        api_key=os.getenv(api_key_env, "").strip() or os.getenv("OCR_API_KEY", "").strip(),
        api_key_env=api_key_env,
        timeout_s=_env_int(_ocr_timeout_env(provider), _env_int("OCR_HTTP_TIMEOUT_S", 120)),
        mode=_ocr_mode(provider),
    )


def _ocr_endpoint_env(provider: str) -> str:
    if provider == "mineru":
        return "MINERU_HTTP_ENDPOINT"
    if provider == "paddleocr":
        return "PADDLEOCR_HTTP_ENDPOINT"
    return "OCR_HTTP_ENDPOINT"


def _ocr_api_key_env(provider: str) -> str:
    if provider == "mineru":
        return "MINERU_API_KEY"
    if provider == "paddleocr":
        return "PADDLEOCR_API_KEY"
    return "OCR_API_KEY"


def _ocr_timeout_env(provider: str) -> str:
    if provider == "mineru":
        return "MINERU_HTTP_TIMEOUT_S"
    if provider == "paddleocr":
        return "PADDLEOCR_HTTP_TIMEOUT_S"
    return "OCR_HTTP_TIMEOUT_S"


def _ocr_mode(provider: str) -> str:
    if provider == "mineru":
        return os.getenv("MINERU_PARSE_MODE", "auto").strip() or "auto"
    if provider == "paddleocr":
        return os.getenv("PADDLEOCR_PIPELINE", "pp_structurev3").strip() or "pp_structurev3"
    return os.getenv("OCR_PARSE_MODE", "http_json").strip() or "http_json"


def _ocr_provider_options(config: OCRProviderConfig) -> dict[str, object]:
    options: dict[str, object] = {
        "return_pages": True,
        "return_blocks": True,
        "return_tables": True,
        "return_layout": True,
    }
    if config.provider == "mineru":
        options["return_markdown"] = True
        options["parse_mode"] = config.mode
    elif config.provider == "paddleocr":
        options["pipeline"] = config.mode
        options["return_confidence"] = True
        options["return_bbox"] = True
    return options


def _ocr_provider_profile(config: OCRProviderConfig) -> dict[str, object]:
    return {
        "provider": config.provider,
        "endpoint_env": config.endpoint_env,
        "api_key_env": config.api_key_env,
        "mode": config.mode,
        "timeout_s": config.timeout_s,
    }


def _parse_docx(content: bytes) -> tuple[str, dict[str, object]]:
    document = DocxDocument(BytesIO(content))
    paragraph_limit = _env_int("KNOWLEDGE_PARSE_MAX_DOCX_PARAGRAPHS", 3000)
    table_row_limit = _env_int("KNOWLEDGE_PARSE_MAX_DOCX_TABLE_ROWS", 3000)
    paragraphs: list[str] = []
    paragraph_count = len(document.paragraphs)
    for paragraph in document.paragraphs[:paragraph_limit]:
        text = paragraph.text.strip()
        if text:
            paragraphs.append(text)
    table_rows: list[str] = []
    table_blocks: list[dict[str, object]] = []
    table_count = len(document.tables)
    parsed_table_rows = 0
    total_table_rows = sum(len(table.rows) for table in document.tables)
    for table_index, table in enumerate(document.tables, start=1):
        current_rows: list[list[str]] = []
        table_truncated = False
        for row_index, row in enumerate(table.rows, start=1):
            if parsed_table_rows >= table_row_limit:
                table_truncated = True
                break
            parsed_table_rows += 1
            values = [cell.text.strip().replace("\n", " ") for cell in row.cells if cell.text.strip()]
            if values:
                current_rows.append(values)
                table_rows.append(f"[Table {table_index} Row {row_index}] " + " | ".join(values))
        if current_rows:
            table_payload: dict[str, object] = {
                "index": table_index,
                "rows": current_rows,
                "extraction": "python-docx",
            }
            if table_truncated:
                table_payload["truncated_after_row_limit"] = True
                table_payload["row_limit"] = table_row_limit
            table_blocks.append(_table_block("docx", table_payload))
        if parsed_table_rows >= table_row_limit:
            break
    return "\n\n".join([*paragraphs, *table_rows]), {
        "docx_paragraph_count": paragraph_count,
        "docx_parsed_paragraph_count": min(paragraph_count, paragraph_limit),
        "docx_paragraph_limit": paragraph_limit,
        "docx_table_count": table_count,
        "docx_table_row_count": total_table_rows,
        "docx_parsed_table_row_count": parsed_table_rows,
        "docx_table_row_limit": table_row_limit,
        "table_blocks": table_blocks,
        "table_block_count": len(table_blocks),
        "table_count": len(table_blocks),
        "truncated_after_parse_limit": paragraph_count > paragraph_limit or total_table_rows > table_row_limit,
    }


def _parse_xlsx(content: bytes) -> tuple[str, dict[str, object]]:
    workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
    formula_workbook = None
    lines: list[str] = []
    max_sheets = _env_int("KNOWLEDGE_PARSE_MAX_XLSX_SHEETS", 20)
    max_rows = _env_int("KNOWLEDGE_PARSE_MAX_XLSX_ROWS_PER_SHEET", 5000)
    max_columns = _env_int("KNOWLEDGE_PARSE_MAX_XLSX_COLUMNS", 80)
    truncated = False
    parsed_sheets = 0
    parsed_rows = 0
    table_blocks: list[dict[str, object]] = []
    try:
        formula_workbook = load_workbook(BytesIO(content), read_only=True, data_only=False)
        sheet_count = len(workbook.worksheets)
        sheets = workbook.worksheets
        formula_sheets = {sheet.title: sheet for sheet in formula_workbook.worksheets}
        for sheet in sheets[:max_sheets]:
            formula_sheet = formula_sheets.get(sheet.title)
            parsed_sheets += 1
            lines.append(f"[Sheet] {sheet.title}")
            sheet_rows: list[list[str]] = []
            data_rows = sheet.iter_rows(max_row=max_rows, max_col=max_columns, values_only=True)
            formula_rows = (
                formula_sheet.iter_rows(max_row=max_rows, max_col=max_columns, values_only=True)
                if formula_sheet
                else []
            )
            for row, formula_row in zip_longest(data_rows, formula_rows, fillvalue=()):
                values = [
                    _cell_text(value, formula_value)
                    for value, formula_value in zip_longest(row or (), formula_row or (), fillvalue=None)
                ]
                values = [value for value in values if value]
                if values:
                    lines.append(" | ".join(values))
                    sheet_rows.append(values)
                    parsed_rows += 1
            formula_max_row = (formula_sheet.max_row or 0) if formula_sheet else 0
            formula_max_column = (formula_sheet.max_column or 0) if formula_sheet else 0
            sheet_truncated = (
                max(sheet.max_row or 0, formula_max_row) > max_rows
                or max(sheet.max_column or 0, formula_max_column) > max_columns
            )
            if (
                sheet_truncated
            ):
                truncated = True
            if sheet_rows:
                table_payload: dict[str, object] = {
                    "index": parsed_sheets,
                    "sheet": sheet.title,
                    "rows": sheet_rows,
                    "extraction": "openpyxl",
                }
                if sheet_truncated:
                    table_payload["truncated_after_row_limit"] = True
                    table_payload["row_limit"] = max_rows
                table_blocks.append(_table_block("xlsx", table_payload))
        if sheet_count > max_sheets:
            truncated = True
    finally:
        workbook.close()
        if formula_workbook is not None:
            formula_workbook.close()
    return "\n".join(lines), {
        "xlsx_sheet_count": sheet_count,
        "xlsx_parsed_sheet_count": parsed_sheets,
        "xlsx_parsed_row_count": parsed_rows,
        "xlsx_sheet_limit": max_sheets,
        "xlsx_row_limit_per_sheet": max_rows,
        "xlsx_column_limit": max_columns,
        "table_blocks": table_blocks,
        "table_block_count": len(table_blocks),
        "table_count": len(table_blocks),
        "truncated_after_parse_limit": truncated,
    }


def _parse_pptx(content: bytes) -> tuple[str, dict[str, object]]:
    presentation = Presentation(BytesIO(content))
    lines: list[str] = []
    max_slides = _env_int("KNOWLEDGE_PARSE_MAX_PPTX_SLIDES", 200)
    max_table_rows = _env_int("KNOWLEDGE_PARSE_MAX_PPTX_TABLE_ROWS", 1000)
    slide_count = len(presentation.slides)
    parsed_slides = 0
    parsed_table_rows = 0
    table_blocks: list[dict[str, object]] = []
    truncated = slide_count > max_slides
    for slide_index, slide in enumerate(presentation.slides, start=1):
        if parsed_slides >= max_slides:
            break
        parsed_slides += 1
        lines.append(f"[Slide {slide_index}]")
        table_index = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
            if getattr(shape, "has_table", False):
                table_index += 1
                table_rows: list[list[str]] = []
                table_truncated = False
                for row in shape.table.rows:
                    if parsed_table_rows >= max_table_rows:
                        truncated = True
                        table_truncated = True
                        break
                    parsed_table_rows += 1
                    values = [cell.text.strip().replace("\n", " ") for cell in row.cells if cell.text.strip()]
                    if values:
                        table_rows.append(values)
                        lines.append(" | ".join(values))
                if table_rows:
                    table_payload: dict[str, object] = {
                        "index": table_index,
                        "slide": slide_index,
                        "rows": table_rows,
                        "extraction": "python-pptx",
                    }
                    if table_truncated:
                        table_payload["truncated_after_row_limit"] = True
                        table_payload["row_limit"] = max_table_rows
                    table_blocks.append(_table_block("pptx", table_payload))
    return "\n".join(lines), {
        "pptx_slide_count": slide_count,
        "pptx_parsed_slide_count": parsed_slides,
        "pptx_slide_limit": max_slides,
        "pptx_parsed_table_row_count": parsed_table_rows,
        "pptx_table_row_limit": max_table_rows,
        "table_blocks": table_blocks,
        "table_block_count": len(table_blocks),
        "table_count": len(table_blocks),
        "truncated_after_parse_limit": truncated,
    }


def _parse_legacy_office(
    payload: KnowledgeProcessRequest,
    content: bytes,
    suffix: str,
) -> tuple[str, int | None, dict[str, object]]:
    target_suffix = LEGACY_OFFICE_TARGETS[suffix]
    conversion = _convert_with_libreoffice(payload.filename, content, target_suffix)
    metadata: dict[str, object] = {"legacy_conversion": conversion}
    if conversion.get("status") != "done":
        metadata["parser"] = "legacy-office-unconverted"
        metadata["needs_human_input"] = True
        return "", None, metadata
    converted = conversion.get("content")
    if not isinstance(converted, bytes):
        metadata["parser"] = "legacy-office-unconverted"
        metadata["needs_human_input"] = True
        return "", None, metadata
    conversion = {key: value for key, value in conversion.items() if key != "content"}
    metadata["legacy_conversion"] = conversion
    if target_suffix == ".docx":
        metadata["parser"] = "libreoffice-python-docx"
        text, parse_metadata = _parse_docx(converted)
        metadata.update(parse_metadata)
        return text, None, metadata
    if target_suffix == ".xlsx":
        metadata["parser"] = "libreoffice-openpyxl"
        text, parse_metadata = _parse_xlsx(converted)
        metadata.update(parse_metadata)
        return text, None, metadata
    if target_suffix == ".pptx":
        metadata["parser"] = "libreoffice-python-pptx"
        text, parse_metadata = _parse_pptx(converted)
        metadata.update(parse_metadata)
        return text, None, metadata
    metadata["parser"] = "legacy-office-unconverted"
    metadata["needs_human_input"] = True
    return "", None, metadata


def _convert_with_libreoffice(filename: str, content: bytes, target_suffix: str) -> dict[str, object]:
    executable = _libreoffice_convert_executable()
    if not executable:
        return {"status": "converter_not_configured", "target_suffix": target_suffix}
    source_name = Path(filename).name or f"input{target_suffix}"
    timeout = _env_int("LIBREOFFICE_CONVERT_TIMEOUT_S", 120)
    with TemporaryDirectory(prefix="zbt-ai-parse-") as tmpdir:
        source_path = Path(tmpdir) / source_name
        source_path.write_bytes(content)
        try:
            completed = subprocess.run(
                [
                    executable,
                    "--headless",
                    "--convert-to",
                    target_suffix.lstrip("."),
                    "--outdir",
                    tmpdir,
                    str(source_path),
                ],
                check=False,
                capture_output=True,
                text=True,
                timeout=timeout,
            )
        except Exception as exc:
            return {
                "status": "failed",
                "target_suffix": target_suffix,
                "error": "conversion_failed",
                "error_type": type(exc).__name__,
            }
        if completed.returncode != 0:
            return {
                "status": "failed",
                "target_suffix": target_suffix,
                "error": "conversion_failed",
                "return_code": completed.returncode,
            }
        converted_files = [
            path
            for path in Path(tmpdir).iterdir()
            if path.suffix.lower() == target_suffix and path.name != source_path.name
        ]
        if not converted_files:
            return {"status": "failed", "target_suffix": target_suffix, "error": "converted file not found"}
        converted_path = max(converted_files, key=lambda path: path.stat().st_mtime)
        return {
            "status": "done",
            "target_suffix": target_suffix,
            "filename": converted_path.name,
            "content": converted_path.read_bytes(),
        }


def _libreoffice_convert_executable() -> str | None:
    return (
        os.getenv("LIBREOFFICE_BIN", "").strip()
        or os.getenv("LIBREOFFICE_PATH", "").strip()
        or shutil.which("soffice")
        or shutil.which("libreoffice")
    )


def _cell_text(value: object, formula_value: object | None = None) -> str:
    if value is None:
        value = formula_value
    if value is None:
        return ""
    return str(value).strip()


def _chunk_text(
    text: str,
    filename: str,
    page_count: int | None,
    max_chunks: int,
) -> list[KnowledgeChunk]:
    normalized = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not normalized:
        return [
            KnowledgeChunk(
                title=filename,
                content=f"{filename} 暂未解析出文本内容，需要人工确认或 OCR。",
                section_path=filename,
                metadata={"needs_human_input": True, "reason": "empty_parse_result"},
            )
        ]

    chunks: list[KnowledgeChunk] = []
    max_chars = 1200
    cursor = 0
    while cursor < len(normalized) and len(chunks) < max_chunks:
        end = min(len(normalized), cursor + max_chars)
        if end < len(normalized):
            boundary = normalized.rfind("\n", cursor, end)
            if boundary > cursor + 300:
                end = boundary
        chunk_text = normalized[cursor:end].strip()
        if chunk_text:
            index = len(chunks) + 1
            page_start = 1 if page_count else None
            page_end = page_count if page_count else None
            chunks.append(
                KnowledgeChunk(
                    title=f"{Path(filename).stem or filename} #{index}",
                    content=chunk_text,
                    section_path=Path(filename).stem or filename,
                    page_start=page_start,
                    page_end=page_end,
                    metadata={"chunk_index": index - 1, "parser_max_chars": max_chars},
                )
            )
        cursor = max(end, cursor + 1)
    if cursor < len(normalized) and chunks:
        chunks[-1].metadata["truncated_after_chunk_limit"] = True
        chunks[-1].metadata["remaining_chars_estimate"] = len(normalized) - cursor
    return chunks


def _summary(text: str, filename: str) -> str:
    compact = " ".join(text.split())
    if not compact:
        return f"{filename} 暂未解析出文本内容，需要人工确认或 OCR。"
    return compact[:180]
